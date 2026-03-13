import streamlit as st
import anthropic
import json
import re
import base64
import io

# ── File parsing imports ──────────────────────────────────────────────────────
try:
    import pdfplumber
except ImportError:
    pdfplumber = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    from docx import Document
except ImportError:
    Document = None

st.set_page_config(page_title="StudyForge", page_icon="✦", layout="centered")

# ── Styling ───────────────────────────────────────────────────────────────────
st.markdown("""
<style>
@import url('https://fonts.googleapis.com/css2?family=Syne:wght@400;600;800&family=DM+Mono:wght@300;400&display=swap');

html, body, [class*="css"] { font-family: 'Syne', sans-serif; }
.stApp { background-color: #0a0a0f; color: #e8e8f0; }
.logo { font-family: 'DM Mono', monospace; font-size: 13px; letter-spacing: 0.2em;
        text-transform: uppercase; color: #c8f060; margin-bottom: 4px; }
h1 { font-size: 2.6rem !important; font-weight: 800 !important;
     letter-spacing: -0.02em !important; line-height: 1.1 !important; }
.tagline { font-family: 'DM Mono', monospace; font-size: 13px;
           color: #6a6a80; margin-top: 4px; margin-bottom: 32px; }
.stTextArea textarea {
    background: #13131a !important; border: 1px solid #2a2a3a !important;
    color: #e8e8f0 !important; font-family: 'DM Mono', monospace !important;
    font-size: 13px !important; border-radius: 10px !important; }
.stButton > button {
    background: #c8f060 !important; color: #0a0a0f !important;
    font-family: 'Syne', sans-serif !important; font-weight: 800 !important;
    font-size: 15px !important; letter-spacing: 0.05em !important;
    border: none !important; border-radius: 10px !important;
    padding: 14px 32px !important; width: 100% !important; }
.stButton > button:hover { background: #d8ff70 !important; }
.card-box { background: #13131a; border: 1px solid #2a2a3a; border-radius: 14px;
            padding: 28px 32px; margin-bottom: 16px; }
.card-box.answer { border-color: #c8f060; background: #1c1c26; }
.badge { font-family: 'DM Mono', monospace; font-size: 10px; letter-spacing: 0.15em;
         text-transform: uppercase; color: #6a6a80; margin-bottom: 10px; }
.badge.q { color: #60a0f0; }
.badge.a { color: #c8f060; }
.card-text { font-size: 17px; font-weight: 600; line-height: 1.5; }
.correct-ans { background: rgba(96,240,160,0.1); border: 1px solid #60f0a0;
               border-radius: 8px; padding: 12px 18px; color: #60f0a0;
               font-weight: 600; margin-bottom: 8px; }
.wrong-ans { background: rgba(240,96,96,0.1); border: 1px solid #f06060;
             border-radius: 8px; padding: 12px 18px; color: #f06060;
             font-weight: 600; margin-bottom: 8px; }
.explanation { background: rgba(200,240,96,0.06); border-left: 3px solid #c8f060;
               border-radius: 0 8px 8px 0; padding: 12px 16px; margin-top: 12px;
               font-family: 'DM Mono', monospace; font-size: 13px;
               color: #9a9ab0; line-height: 1.6; }
.score-box { background: #13131a; border: 1px solid #c8f060; border-radius: 14px;
             padding: 40px; text-align: center; margin-top: 32px; }
.score-big { font-size: 72px; font-weight: 800; color: #c8f060; line-height: 1; }
.score-sub { font-family: 'DM Mono', monospace; font-size: 13px;
             color: #6a6a80; margin-top: 8px; }
.divider { border: none; border-top: 1px solid #2a2a3a; margin: 32px 0; }
.file-info { background: #13131a; border: 1px solid #2a2a3a; border-radius: 10px;
             padding: 14px 18px; font-family: 'DM Mono', monospace; font-size: 12px;
             color: #9a9ab0; margin-bottom: 16px; }
.file-info strong { color: #c8f060; }
</style>
""", unsafe_allow_html=True)


# ── API client ────────────────────────────────────────────────────────────────
@st.cache_resource
def get_client():
    return anthropic.Anthropic(api_key=st.secrets["ANTHROPIC_API_KEY"])


# ── File extraction ───────────────────────────────────────────────────────────
def extract_pdf(file_bytes: bytes) -> str:
    with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
        pages = [page.extract_text() or "" for page in pdf.pages]
    return "\n\n".join(pages)


def extract_pptx(file_bytes: bytes) -> str:
    prs = Presentation(io.BytesIO(file_bytes))
    slides_text = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        if texts:
            slides_text.append(f"[Slide {i}]\n" + "\n".join(texts))
    return "\n\n".join(slides_text)


def extract_docx(file_bytes: bytes) -> str:
    doc = Document(io.BytesIO(file_bytes))
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def extract_image(file_bytes: bytes, mime_type: str) -> str:
    client = get_client()
    b64 = base64.standard_b64encode(file_bytes).decode("utf-8")
    message = client.messages.create(
        model="claude-opus-4-5",
        max_tokens=2000,
        messages=[{
            "role": "user",
            "content": [
                {
                    "type": "image",
                    "source": {"type": "base64", "media_type": mime_type, "data": b64}
                },
                {
                    "type": "text",
                    "text": "This is a lecture slide or study material image. Extract and transcribe ALL text visible in the image, preserving structure where possible. Include headings, bullet points, labels, and any other text. Output only the extracted text."
                }
            ]
        }]
    )
    return message.content[0].text.strip()


def process_uploaded_file(uploaded_file) -> tuple[str, str]:
    file_bytes = uploaded_file.read()
    name = uploaded_file.name.lower()

    if name.endswith(".pdf"):
        text = extract_pdf(file_bytes)
        return text, f"📄 Extracted from **{uploaded_file.name}**"
    elif name.endswith(".pptx"):
        text = extract_pptx(file_bytes)
        return text, f"📊 Extracted from **{uploaded_file.name}**"
    elif name.endswith(".docx"):
        text = extract_docx(file_bytes)
        return text, f"📝 Extracted from **{uploaded_file.name}**"
    elif name.endswith((".jpg", ".jpeg")):
        text = extract_image(file_bytes, "image/jpeg")
        return text, f"🖼️ Extracted from **{uploaded_file.name}** via vision"
    elif name.endswith(".png"):
        text = extract_image(file_bytes, "image/png")
        return text, f"🖼️ Extracted from **{uploaded_file.name}** via vision"
    else:
        return "", "Unsupported file type."


# ── Study set generation ──────────────────────────────────────────────────────
def generate_study_set(notes: str, fc_count: int, quiz_count: int) -> dict:
    client = get_client()
    prompt = f"""You are a study assistant. Given the following lecture notes, generate exactly {fc_count} flashcards and exactly {quiz_count} multiple-choice quiz questions.

Return ONLY a JSON object in this exact format with NO other text:
{{
  "flashcards": [
    {{ "question": "...", "answer": "..." }}
  ],
  "quiz": [
    {{
      "question": "...",
      "options": ["A: ...", "B: ...", "C: ...", "D: ..."],
      "correct": 0,
      "explanation": "..."
    }}
  ]
}}

Rules:
- Flashcards should cover key concepts, definitions, and facts
- Quiz questions should be clear multiple choice with 4 options (correct is index 0-3)
- Keep questions concise and academically rigorous
- Explanations should clarify why the answer is correct
- Return ONLY valid JSON, no markdown, no backticks

LECTURE NOTES:
{notes[:8000]}"""

    message = client.messages.create(
        model="claude-opus-4-5",
        max_tokens=4000,
        messages=[{"role": "user", "content": prompt}]
    )

    raw = message.content[0].text.strip()
    clean = re.sub(r"```json|```", "", raw).strip()
    return json.loads(clean)


# ── Session state defaults ────────────────────────────────────────────────────
for key, default in [
    ("study_set", None),
    ("fc_index", 0),
    ("fc_flipped", False),
    ("quiz_answers", {}),
    ("quiz_done", False),
    ("view", "upload"),
    ("extracted_text", ""),
    ("file_status", ""),
]:
    if key not in st.session_state:
        st.session_state[key] = default


# ── HEADER ────────────────────────────────────────────────────────────────────
st.markdown('<div class="logo">✦ StudyForge</div>', unsafe_allow_html=True)
st.markdown("# Turn your notes into a **study arsenal.**")
st.markdown('<div class="tagline">// upload slides or paste notes → flashcards + quiz in seconds</div>', unsafe_allow_html=True)
st.markdown('<hr class="divider">', unsafe_allow_html=True)


# ══════════════════════════════════════════════════════════════════════════════
# UPLOAD VIEW
# ══════════════════════════════════════════════════════════════════════════════
if st.session_state.view == "upload":

    uploaded_file = st.file_uploader(
        "Upload lecture slides or notes",
        type=["pdf", "pptx", "docx", "jpg", "jpeg", "png"],
        help="Supported: PDF, PowerPoint (.pptx), Word (.docx), JPG, PNG"
    )

    if uploaded_file is not None:
        if st.button("Extract text from file"):
            with st.spinner("Extracting content..."):
                try:
                    text, status = process_uploaded_file(uploaded_file)
                    if text:
                        st.session_state.extracted_text = text
                        st.session_state.file_status = status
                        st.rerun()
                    else:
                        st.error("Could not extract text from this file.")
                except Exception as e:
                    st.error(f"Extraction failed: {e}")

    if st.session_state.file_status:
        word_count = len(st.session_state.extracted_text.split())
        st.markdown(
            f'<div class="file-info">{st.session_state.file_status} — '
            f'<strong>{word_count} words</strong> extracted</div>',
            unsafe_allow_html=True
        )

    st.markdown("**Or paste notes directly:**")
    notes_input = st.text_area(
        "Notes",
        value=st.session_state.extracted_text,
        height=240,
        placeholder="Paste your lecture notes here, or upload a file above...",
        label_visibility="collapsed",
        key="notes_textarea"
    )

    col1, col2 = st.columns(2)
    with col1:
        fc_count = st.select_slider("Flashcards", options=[5, 10, 15], value=10)
    with col2:
        quiz_count = st.select_slider("Quiz questions", options=[5, 8, 12], value=8)

    st.markdown("<br>", unsafe_allow_html=True)

    if st.button("Generate Study Set →"):
        notes = notes_input.strip()
        if not notes or len(notes) < 50:
            st.error("Please add some content first — upload a file or paste notes.")
        else:
            with st.spinner("Generating your study set..."):
                try:
                    result = generate_study_set(notes, fc_count, quiz_count)
                    st.session_state.study_set = result
                    st.session_state.fc_index = 0
                    st.session_state.fc_flipped = False
                    st.session_state.quiz_answers = {}
                    st.session_state.quiz_done = False
                    st.session_state.view = "results"
                    st.rerun()
                except Exception as e:
                    st.error(f"Something went wrong: {e}")


# ══════════════════════════════════════════════════════════════════════════════
# RESULTS VIEW
# ══════════════════════════════════════════════════════════════════════════════
else:
    data = st.session_state.study_set
    flashcards = data.get("flashcards", [])
    quiz = data.get("quiz", [])

    c1, c2, c3 = st.columns(3)
    c1.metric("Flashcards", len(flashcards))
    c2.metric("Quiz questions", len(quiz))
    c3.metric("Quiz progress", f"{len(st.session_state.quiz_answers)}/{len(quiz)}")

    if st.button("← New notes"):
        st.session_state.view = "upload"
        st.session_state.study_set = None
        st.session_state.extracted_text = ""
        st.session_state.file_status = ""
        st.rerun()

    st.markdown('<hr class="divider">', unsafe_allow_html=True)

    tab1, tab2 = st.tabs(["✦ Flashcards", "✎ Quiz"])

    # ── FLASHCARDS ────────────────────────────────────────────────────────────
    with tab1:
        if not flashcards:
            st.warning("No flashcards generated.")
        else:
            idx = st.session_state.fc_index
            card = flashcards[idx]

            st.markdown(f"**Card {idx + 1} of {len(flashcards)}**")
            st.progress((idx + 1) / len(flashcards))
            st.markdown("<br>", unsafe_allow_html=True)

            st.markdown(f"""
            <div class="card-box">
                <div class="badge q">Question</div>
                <div class="card-text">{card['question']}</div>
            </div>""", unsafe_allow_html=True)

            flip_label = "Hide answer ↑" if st.session_state.fc_flipped else "Reveal answer ↓"
            if st.button(flip_label, key="flip"):
                st.session_state.fc_flipped = not st.session_state.fc_flipped
                st.rerun()

            if st.session_state.fc_flipped:
                st.markdown(f"""
                <div class="card-box answer">
                    <div class="badge a">Answer</div>
                    <div class="card-text">{card['answer']}</div>
                </div>""", unsafe_allow_html=True)

            st.markdown("<br>", unsafe_allow_html=True)
            prev_col, _, next_col = st.columns([1, 2, 1])
            with prev_col:
                if st.button("← Prev", disabled=idx == 0, key="prev"):
                    st.session_state.fc_index -= 1
                    st.session_state.fc_flipped = False
                    st.rerun()
            with next_col:
                if st.button("Next →", disabled=idx == len(flashcards) - 1, key="next"):
                    st.session_state.fc_index += 1
                    st.session_state.fc_flipped = False
                    st.rerun()

    # ── QUIZ ──────────────────────────────────────────────────────────────────
    with tab2:
        if not quiz:
            st.warning("No quiz questions generated.")
        else:
            for qi, q in enumerate(quiz):
                st.markdown(f"**Q{qi + 1}.** {q['question']}")
                answered_this = qi in st.session_state.quiz_answers

                if not answered_this:
                    for oi, opt in enumerate(q["options"]):
                        if st.button(opt, key=f"q{qi}_o{oi}"):
                            st.session_state.quiz_answers[qi] = oi
                            if len(st.session_state.quiz_answers) == len(quiz):
                                st.session_state.quiz_done = True
                            st.rerun()
                else:
                    chosen = st.session_state.quiz_answers[qi]
                    for oi, opt in enumerate(q["options"]):
                        if oi == q["correct"]:
                            st.markdown(f'<div class="correct-ans">✓ {opt}</div>', unsafe_allow_html=True)
                        elif oi == chosen and chosen != q["correct"]:
                            st.markdown(f'<div class="wrong-ans">✗ {opt}</div>', unsafe_allow_html=True)
                        else:
                            st.markdown(f"<div style='padding:10px 18px;color:#6a6a80'>{opt}</div>", unsafe_allow_html=True)
                    if q.get("explanation"):
                        st.markdown(f'<div class="explanation">💡 {q["explanation"]}</div>', unsafe_allow_html=True)

                st.markdown("<br>", unsafe_allow_html=True)

            if st.session_state.quiz_done:
                correct = sum(
                    1 for qi, chosen in st.session_state.quiz_answers.items()
                    if chosen == quiz[qi]["correct"]
                )
                pct = round((correct / len(quiz)) * 100)
                st.markdown(f"""
                <div class="score-box">
                    <div class="score-big">{pct}%</div>
                    <div class="score-sub">{correct} / {len(quiz)} correct</div>
                </div>""", unsafe_allow_html=True)
                st.markdown("<br>", unsafe_allow_html=True)
                if st.button("Retry Quiz"):
                    st.session_state.quiz_answers = {}
                    st.session_state.quiz_done = False
                    st.rerun()
